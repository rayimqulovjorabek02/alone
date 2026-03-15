"""
backend/app/core/file_reader.py — Fayl o'qish yordamchisi
PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, JSON, PY, JS va boshqalar
"""
import os


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt",
    ".txt", ".md", ".csv", ".json", ".py", ".js", ".ts",
    ".html", ".xml", ".yaml", ".yml", ".log",
}


def read_file(path: str) -> str:
    """Fayl turini aniqlab matnini o'qish."""
    ext = os.path.splitext(path)[-1].lower()

    readers = {
        ".pdf":  _read_pdf,
        ".docx": _read_docx,
        ".doc":  _read_docx,
        ".xlsx": _read_excel,
        ".xls":  _read_excel,
        ".pptx": _read_pptx,
        ".ppt":  _read_pptx,
        ".csv":  _read_csv,
    }

    reader = readers.get(ext)
    if reader:
        return reader(path)
    return _read_text(path)


def _read_pdf(path: str) -> str:
    # pypdf (yangi)
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        pages = [p.extract_text() or "" for p in reader.pages]
        text = "\n".join(pages).strip()
        if text:
            return text
    except Exception as e:
        print(f"[pypdf xato]: {e}")

    # PyPDF2 fallback
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = [p.extract_text() or "" for p in reader.pages]
            return "\n".join(pages).strip()
    except Exception as e:
        print(f"[PyPDF2 xato]: {e}")

    return "[PDF o'qib bo'lmadi]"


def _read_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []

        # Paragraflar
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)

        # Jadvallar
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        return "\n".join(parts)
    except ImportError:
        return "[python-docx o'rnatilmagan: pip install python-docx]"
    except Exception as e:
        return f"[DOCX o'qish xato: {e}]"


def _read_excel(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"--- {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) for cell in row if cell is not None
                )
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        pass

    # xlrd fallback (eski .xls)
    try:
        import xlrd
        wb = xlrd.open_workbook(path)
        parts = []
        for sheet in wb.sheets():
            parts.append(f"--- {sheet.name} ---")
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                row_text = " | ".join(str(c) for c in row if str(c).strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        return "[openpyxl o'rnatilmagan: pip install openpyxl]"
    except Exception as e:
        return f"[Excel o'qish xato: {e}]"


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slayd {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                # Jadval
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        return "[python-pptx o'rnatilmagan: pip install python-pptx]"
    except Exception as e:
        return f"[PPTX o'qish xato: {e}]"


def _read_csv(path: str) -> str:
    try:
        import csv
        parts = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[CSV o'qish xato: {e}]"


def _read_text(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[Fayl o'qish xato: {e}]"


def get_file_info(path: str) -> dict:
    """Fayl haqida ma'lumot."""
    stat = os.stat(path)
    ext  = os.path.splitext(path)[-1].lower()
    return {
        "name":      os.path.basename(path),
        "size_kb":   round(stat.st_size / 1024, 1),
        "extension": ext,
        "supported": ext in SUPPORTED_EXTENSIONS,
    }